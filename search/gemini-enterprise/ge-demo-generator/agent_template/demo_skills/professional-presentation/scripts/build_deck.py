# Copyright 2026 Google LLC
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     https://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.


# Deployed as a runtime template into the user's Cloud Shell (not imported by
# repo tooling); validated by py_compile and end-to-end demo deployments.
# Repo-level strict lint/typing is intentionally skipped for this generated-
# origin runtime code; incremental typing is planned as follow-up.
# flake8: noqa
# pylint: skip-file
# mypy: ignore-errors
# ruff: noqa

# Scaffold for a professional 16:9 deck. Copy and adapt; do not import this file.
# Usage pattern: define SLIDES as data, then render mechanically. Keeping content
# as data makes the verification pass and later edits trivial.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRIMARY = RGBColor(0x1A, 0x2B, 0x4A)
ACCENT = RGBColor(0x2E, 0x6F, 0xDB)
WARN = RGBColor(0xE8, 0xA3, 0x3D)
BODY = RGBColor(0x33, 0x3F, 0x50)
MUTED = RGBColor(0x8A, 0x93, 0xA6)
PANEL = RGBColor(0xF2, 0xF5, 0xFA)

MARGIN = Inches(0.6)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # fully blank layout


def text_box(slide, left, top, width, height, text, size, color=BODY,
             bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return box


def accent_rule(slide, top=Inches(1.15)):
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, top,
                                 SLIDE_W - MARGIN * 2, Emu(27432))  # 0.03in
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    return bar


def content_slide(prs, title, bullets, footer_text, page_no):
    slide = blank(prs)
    text_box(slide, MARGIN, Inches(0.35), SLIDE_W - MARGIN * 2, Inches(0.8),
             title, 30, PRIMARY, bold=True)
    accent_rule(slide)
    box = slide.shapes.add_textbox(MARGIN, Inches(1.5),
                                   SLIDE_W - MARGIN * 2, Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.size = Pt(16)
        run.font.color.rgb = BODY
        run.font.name = "Arial"
        p.space_after = Pt(10)
    text_box(slide, MARGIN, Inches(7.05), Inches(6), Inches(0.3),
             footer_text, 10, MUTED)
    text_box(slide, SLIDE_W - MARGIN - Inches(1), Inches(7.05), Inches(1),
             Inches(0.3), str(page_no), 10, MUTED, align=PP_ALIGN.RIGHT)
    return slide


def verify(path, expected_slides):
    prs = Presentation(path)
    assert len(prs.slides) == expected_slides, (
        "slide count %d != expected %d" % (len(prs.slides), expected_slides))
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            assert shape.left >= 0 and shape.top >= 0, (
                "slide %d: shape off-canvas" % idx)
            assert shape.left + shape.width <= SLIDE_W + Emu(1), (
                "slide %d: shape exceeds right edge" % idx)
            assert shape.top + shape.height <= SLIDE_H + Emu(1), (
                "slide %d: shape exceeds bottom edge" % idx)
    print("verify OK:", path)
